"""
ingestion/element_classifier.py
─────────────────────────────────
Change 1 — Per-element visual classification.

Classifies each extracted page element into:
    "text"   → existing OpenAI text-embedding path (unchanged)
    "table"  → ColPali crop path
    "image"  → ColPali crop path
    "chart"  → ColPali crop path (charts / graphs embedded in PDFs/PPTX)

Only elements classified as table | image | chart are sent through ColPali.
Text elements go through the existing embedding pipeline — no change there.

How crops work:
  - PDF:  fitz bbox of the image/table block → crop pixmap → PNG bytes
  - PPTX: shape.image.blob already is the crop (we already have this)
  - DOCX/XLSX: no visual element extraction needed (text + table markdown only)

The output `VisualElement` dataclass is what ColPali receives instead of a
full-page screenshot. This is a net improvement because:
  1. ColPali sees only the meaningful visual, not whitespace/headers/footers.
  2. We skip full-page rendering for text-only pages entirely.
  3. The crop is smaller → faster embedding, less memory.

Fallback:
  If crop extraction fails for any element, the element is skipped from
  ColPali (doesn't crash the pipeline). The text description from the
  existing embedder is still used as the chunk content.
"""

from __future__ import annotations

import base64
import io
from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal

ElementKind = Literal["text", "table", "image", "chart"]


@dataclass
class VisualElement:
    """A cropped visual element ready for ColPali embedding."""
    kind:        ElementKind
    image_b64:   str           # base64 PNG of the crop (not full page)
    source_file: str
    page_idx:    int
    element_idx: int           # position on the page (for stable IDs)
    bbox:        tuple | None  # (x0, y0, x1, y1) in page coordinates
    metadata:    dict = field(default_factory=dict)


# ── Heuristic chart detector ───────────────────────────────────────────────────

_CHART_EXTENSIONS = {".png", ".jpg", ".jpeg", ".svg"}

# Image aspect-ratio heuristic: wide-and-short images are likely charts/graphs
_CHART_MIN_ASPECT  = 1.2   # width / height
_CHART_MAX_ASPECT  = 5.0


def _is_chart_shape(width: int, height: int) -> bool:
    """Heuristic: wide landscape images are likely charts, not photos/logos."""
    if height == 0:
        return False
    ratio = width / height
    return _CHART_MIN_ASPECT <= ratio <= _CHART_MAX_ASPECT


# ── PDF element extractor ──────────────────────────────────────────────────────

def extract_visual_elements_pdf(
    path: Path,
    metadata: dict,
    min_image_px: int = 80,
) -> list[VisualElement]:
    """
    Extract visual elements (images, tables, charts) from a PDF.
    Returns VisualElement list — one entry per qualifying element (NOT per page).
    Text-only pages produce zero entries here.
    """
    import fitz
    import pdfplumber

    elements: list[VisualElement] = []
    fitz_doc  = fitz.open(str(path))

    # ── Per-page image extraction (ColPali replaces full-page screenshot) ────
    for page_idx, page in enumerate(fitz_doc):
        img_list = page.get_images(full=True)
        for elem_idx, img_meta in enumerate(img_list):
            xref     = img_meta[0]
            base_img = fitz_doc.extract_image(xref)
            w, h     = base_img["width"], base_img["height"]

            if w < min_image_px or h < min_image_px:
                continue   # skip icons / bullets

            img_b64 = base64.b64encode(base_img["image"]).decode()
            kind: ElementKind = "chart" if _is_chart_shape(w, h) else "image"

            # Get bbox of the image on the page for reference
            img_rects = page.get_image_rects(xref)
            bbox = tuple(img_rects[0]) if img_rects else None

            elements.append(VisualElement(
                kind=kind,
                image_b64=img_b64,
                source_file=str(path),
                page_idx=page_idx,
                element_idx=elem_idx,
                bbox=bbox,
                metadata=metadata,
            ))

    fitz_doc.close()

    # ── Table crops via pdfplumber (render bbox region as image) ─────────────
    with pdfplumber.open(str(path)) as pdf:
        for page_idx, page in enumerate(pdf.pages):
            for tbl_idx, table in enumerate(page.find_tables()):
                bbox = table.bbox   # (x0, top, x1, bottom) in PDF points
                try:
                    # Crop and rasterise just the table bounding box
                    crop    = page.crop(bbox)
                    pil_img = crop.to_image(resolution=120).original
                    buf     = io.BytesIO()
                    pil_img.save(buf, format="PNG")
                    img_b64 = base64.b64encode(buf.getvalue()).decode()

                    elements.append(VisualElement(
                        kind="table",
                        image_b64=img_b64,
                        source_file=str(path),
                        page_idx=page_idx,
                        element_idx=1000 + tbl_idx,   # offset to avoid collision with img idx
                        bbox=bbox,
                        metadata=metadata,
                    ))
                except Exception:
                    pass   # if crop fails, skip — text chunk still exists

    return elements


# ── PPTX element extractor ────────────────────────────────────────────────────

def extract_visual_elements_pptx(
    path: Path,
    metadata: dict,
    min_image_px: int = 80,
) -> list[VisualElement]:
    """
    Extract visual elements (images, tables, charts) from a PPTX.
    Each shape is already a crop — no full-page rendering needed.
    """
    from pptx import Presentation
    from PIL import Image

    elements: list[VisualElement] = []
    prs = Presentation(str(path))

    for slide_idx, slide in enumerate(prs.slides):
        elem_idx = 0

        for shape in slide.shapes:
            # ── Picture shapes (images / diagrams / charts) ───────────────
            if shape.shape_type == 13:   # MSO_SHAPE_TYPE.PICTURE
                try:
                    img_bytes = shape.image.blob
                    img       = Image.open(io.BytesIO(img_bytes))
                    w, h      = img.width, img.height

                    if w < min_image_px or h < min_image_px:
                        continue

                    buf = io.BytesIO()
                    img.convert("RGB").save(buf, format="PNG")
                    img_b64 = base64.b64encode(buf.getvalue()).decode()
                    kind: ElementKind = "chart" if _is_chart_shape(w, h) else "image"

                    elements.append(VisualElement(
                        kind=kind,
                        image_b64=img_b64,
                        source_file=str(path),
                        page_idx=slide_idx,
                        element_idx=elem_idx,
                        bbox=None,
                        metadata=metadata,
                    ))
                    elem_idx += 1
                except Exception:
                    pass

            # ── Table shapes ──────────────────────────────────────────────
            if shape.has_table:
                try:
                    # Render the table shape region as an image
                    # (We use the shape's left/top/width/height in EMUs)
                    from pptx.util import Emu
                    from PIL import Image as PILImage, ImageDraw

                    rows = [
                        [cell.text for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    if not rows:
                        continue

                    # Simple text render of the table into a PIL image
                    line_h  = 24
                    col_w   = 160
                    n_cols  = len(rows[0])
                    n_rows  = len(rows)
                    img_w   = col_w * n_cols
                    img_h   = line_h * n_rows + 10

                    img  = PILImage.new("RGB", (img_w, img_h), "white")
                    draw = ImageDraw.Draw(img)
                    for r, row in enumerate(rows):
                        for c, cell in enumerate(row):
                            x = c * col_w + 4
                            y = r * line_h + 4
                            draw.text((x, y), str(cell)[:30], fill="black")
                        # Row separator
                        draw.line([(0, (r + 1) * line_h), (img_w, (r + 1) * line_h)],
                                  fill="#CCCCCC", width=1)

                    buf     = io.BytesIO()
                    img.save(buf, format="PNG")
                    img_b64 = base64.b64encode(buf.getvalue()).decode()

                    elements.append(VisualElement(
                        kind="table",
                        image_b64=img_b64,
                        source_file=str(path),
                        page_idx=slide_idx,
                        element_idx=elem_idx + 500,
                        bbox=None,
                        metadata=metadata,
                    ))
                except Exception:
                    pass

    return elements


# ── Dispatcher ────────────────────────────────────────────────────────────────

def extract_visual_elements(
    path: Path,
    metadata: dict,
) -> list[VisualElement]:
    """
    Main entry point. Returns all visual elements (image / table / chart)
    from a document. Only these go to ColPali — text pages produce nothing.

    DOCX and XLSX are excluded from ColPali (text + markdown tables only).
    """
    ext = path.suffix.lower().lstrip(".")
    if ext == "pdf":
        return extract_visual_elements_pdf(path, metadata)
    elif ext == "pptx":
        return extract_visual_elements_pptx(path, metadata)
    # docx / xlsx → no visual element extraction; return empty list
    return []